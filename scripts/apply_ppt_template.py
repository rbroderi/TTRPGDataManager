"""Apply a PowerPoint template to a D2-generated deck.

This script copies every slide from the source deck (typically the output of
`d2 input.d2 output.pptx`) into a template presentation while preserving the
slide master, theme, and background assets.
"""

from __future__ import annotations

import argparse
from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import Any

import pptx
import pptx.enum.shapes
import pptx.oxml.ns


def _clear_template_slides(pres: Any) -> None:
    """Remove existing slides from the template presentation."""
    slide_ids = list(pres.slides._sldIdLst)  # noqa: SLF001
    for slide_id in reversed(slide_ids):
        r_id = slide_id.rId
        pres.part.drop_rel(r_id)
        pres.slides._sldIdLst.remove(slide_id)  # noqa: SLF001


def _clone_shape(src_shape: Any, dest_slide: Any) -> None:
    """Clone *src_shape* onto *dest_slide*."""
    new_shape_element = deepcopy(src_shape._element)  # noqa: SLF001

    if src_shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
        blip = new_shape_element.xpath(".//a:blip")
        if blip:
            old_r_id = blip[0].get(pptx.oxml.ns.qn("r:embed"))
            if old_r_id:
                part = src_shape.part
                if hasattr(part, "related_part"):
                    image_part = part.related_part(old_r_id)
                else:  # pragma: no cover - older python-pptx fallback
                    image_part = part.related_parts[old_r_id]

                image_stream = BytesIO(image_part.blob)
                _, new_r_id = dest_slide.part.get_or_add_image_part(image_stream)
                blip[0].set(pptx.oxml.ns.qn("r:embed"), new_r_id)

    dest_slide.shapes._spTree.insert_element_before(  # noqa: SLF001
        new_shape_element,
        "p:extLst",
    )


def _copy_slide(
    src_slide: Any,
    dest_presentation: Any,
    layout_name: str | None,
) -> None:
    """Copy *src_slide* into *dest_presentation*."""
    if layout_name:
        layout = next(
            (
                layout
                for layout in dest_presentation.slide_layouts
                if layout.name == layout_name
            ),
            None,
        )
    else:
        layout = None

    layout = layout or dest_presentation.slide_layouts[0]
    new_slide = dest_presentation.slides.add_slide(layout)

    for shape in src_slide.shapes:
        _clone_shape(shape, new_slide)


def apply_template(
    source: Path,
    template: Path,
    output: Path,
    layout_name: str | None,
) -> None:
    """Copy all slides from *source* into *template* and write *output*."""
    src_pres = pptx.Presentation(str(source.resolve()))
    template_pres = pptx.Presentation(str(template.resolve()))

    _clear_template_slides(template_pres)

    for slide in src_pres.slides:
        _copy_slide(slide, template_pres, layout_name)

    template_pres.save(str(output.resolve()))


def parse_args() -> argparse.Namespace:
    """Parse command-line arguments."""
    parser = argparse.ArgumentParser(description="Apply PPT template to D2 deck.")
    parser.add_argument(
        "--source",
        "-s",
        type=Path,
        required=True,
        help="Path to the D2-generated PPTX file.",
    )
    parser.add_argument(
        "--template",
        "-t",
        type=Path,
        required=True,
        help="Path to the template PPTX file.",
    )
    parser.add_argument(
        "--output",
        "-o",
        type=Path,
        required=True,
        help="Destination PPTX path.",
    )
    parser.add_argument(
        "--layout",
        type=str,
        default=None,
        help="Optional slide layout name to use (defaults to first layout).",
    )
    return parser.parse_args()


def main() -> None:
    """Entrypoint for CLI usage."""
    args = parse_args()
    apply_template(args.source, args.template, args.output, args.layout)


if __name__ == "__main__":
    main()
